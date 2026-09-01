from __future__ import annotations

import argparse
import html
import os
import re
import shutil
import sys
import threading
import webbrowser
from collections import defaultdict
from pathlib import Path
from typing import Callable
from xml.etree import ElementTree as ET

try:
    from openpyxl import Workbook, load_workbook
    from openpyxl.styles import Alignment, Font, PatternFill
    from openpyxl.utils import get_column_letter
    from openpyxl.utils.cell import range_boundaries
    from openpyxl.worksheet.table import Table, TableStyleInfo
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Pt
except ModuleNotFoundError as exc:
    missing = exc.name or "a required package"
    print(
        f"Missing Python dependency: {missing}.\n"
        "Use a corporate-approved Python environment containing the packages "
        "listed in requirements.txt: python-pptx, openpyxl and Pillow.",
        file=sys.stderr,
    )
    raise SystemExit(2) from exc


ROOT = Path(__file__).resolve().parent
MODELS_DIR = ROOT / "models"
OUTPUT_DIR = ROOT / "output"
PREVIEW_DIR = ROOT / "preview"
DEFAULT_MODEL = "usd_cash_allocation"
PX_TO_PT = 0.75
DRAWING_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

TABLE_SCHEMAS = {
    "Control": (
        "tblSettings",
        ["Setting", "Value", "Purpose"],
    ),
    "Formats": (
        "tblFormats",
        [
            "Format_ID", "Role", "Font_Family", "Font_Size_pt", "Bold", "Italic",
            "Font_Color_Light_Page", "Font_Color_Dark_Page", "Alignment",
            "Vertical_Alignment", "Usage",
        ],
    ),
    "Slides": (
        "tblSlides",
        ["Slide_No", "Slide_ID", "Width_px", "Height_px", "Background_Color", "Theme_Mode", "Layout_Name"],
    ),
    "Objects": (
        "tblObjects",
        [
            "Enabled", "Slide_No", "Object_ID", "Object_Name", "Z_Order", "Object_Type", "Geometry",
            "Text_Content", "Asset_ID", "Asset_Row", "Left_px", "Top_px", "Width_px", "Height_px",
            "Rotation_deg", "Format_ID", "Fill_Color", "Fill_Transparency", "Line_Color",
            "Line_Transparency", "Line_Width_px", "Text_Alignment", "Vertical_Alignment",
            "Margin_Left_px", "Margin_Right_px", "Margin_Top_px", "Margin_Bottom_px",
        ],
    ),
    "Text Runs": (
        "tblTextRuns",
        [
            "Object_ID", "Slide_No", "Paragraph_No", "Run_No", "Start_Char", "Character_Count",
            "Run_Text", "Font_Family", "Font_Size_pt", "Font_Color", "Bold", "Italic",
        ],
    ),
    "Paragraphs": (
        "tblParagraphs",
        ["Object_ID", "Slide_No", "Paragraph_No", "Start_Char", "Character_Count", "Alignment", "Line_Spacing_Ratio"],
    ),
    "Assets": (
        "tblAssets",
        ["Asset_No", "Asset_ID", "Asset_Row", "Description"],
    ),
}


def safe_model_name(value: str) -> str:
    clean = re.sub(r"[^a-zA-Z0-9_-]+", "_", value.strip()).strip("_").lower()
    if not clean:
        raise ValueError("Model name cannot be blank.")
    return clean


def model_dir(model_name: str) -> Path:
    return MODELS_DIR / safe_model_name(model_name)


def workbook_path(model_name: str) -> Path:
    return model_dir(model_name) / "model.xlsx"


def asset_dir(model_name: str) -> Path:
    return model_dir(model_name) / "assets"


def reference_dir(model_name: str) -> Path:
    return model_dir(model_name) / "reference"


def ensure_root_structure() -> None:
    MODELS_DIR.mkdir(parents=True, exist_ok=True)
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    PREVIEW_DIR.mkdir(parents=True, exist_ok=True)


def list_models() -> list[str]:
    ensure_root_structure()
    return sorted(path.name for path in MODELS_DIR.iterdir() if path.is_dir() and (path / "model.xlsx").exists())


def _write_table_sheet(workbook: Workbook, sheet_name: str, table_name: str, headers: list[str], rows: list[list]) -> None:
    sheet = workbook.create_sheet(sheet_name)
    sheet.append(headers)
    for row in rows or [[None] * len(headers)]:
        sheet.append(row)
    end_row = max(2, sheet.max_row)
    end_col = len(headers)
    table = Table(displayName=table_name, ref=f"A1:{get_column_letter(end_col)}{end_row}")
    table.tableStyleInfo = TableStyleInfo(
        name="TableStyleMedium2",
        showFirstColumn=False,
        showLastColumn=False,
        showRowStripes=True,
        showColumnStripes=False,
    )
    sheet.add_table(table)
    sheet.freeze_panes = "A2"
    sheet.row_dimensions[1].height = 24
    for cell in sheet[1]:
        cell.font = Font(name="Arial", size=10, bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor="2F6BFF")
        cell.alignment = Alignment(vertical="center")
    for column_index, header in enumerate(headers, start=1):
        width = min(max(len(str(header)) + 3, 12), 34)
        sheet.column_dimensions[get_column_letter(column_index)].width = width


def create_model(model_name: str) -> Path:
    name = safe_model_name(model_name)
    destination = model_dir(name)
    if workbook_path(name).exists():
        raise FileExistsError(f"Model already exists: {name}")
    destination.mkdir(parents=True, exist_ok=True)
    asset_dir(name).mkdir(exist_ok=True)
    reference_dir(name).mkdir(exist_ok=True)

    workbook = Workbook()
    workbook.remove(workbook.active)
    settings_rows = [
        ["OUTPUT_FILE_NAME", f"{name}.pptx", "PowerPoint output filename"],
        ["SLIDE_WIDTH_PX", 768, "Portrait canvas width"],
        ["SLIDE_HEIGHT_PX", 960, "Portrait canvas height"],
        ["PX_TO_PT", 0.75, "Pixel-to-point conversion"],
        ["USE_EXACT_RUN_FORMATTING", False, "Use Text Runs instead of simplified Formats"],
    ]
    format_rows = [
        ["TITLE", "Title", "Arial", 38, False, False, "#0B1826", "#FFFFFF", "left", "top", "Main slide title"],
        ["SUBTITLE", "Subtitle", "Arial", 18, False, False, "#526477", "#C5D1DC", "left", "top", "Supporting message"],
        ["SECTION", "Section", "Arial", 12, True, False, "#2F6BFF", "#4B8CFF", "left", "top", "Section label"],
        ["BODY", "Body", "Arial", 16, False, False, "#0B1826", "#FFFFFF", "left", "top", "Body copy"],
        ["FOOTER", "Footer", "Arial", 8, False, False, "#526477", "#9EB0C0", "left", "bottom", "Footer or source"],
    ]
    slide_rows = [[1, "S01", 768, 960, "#FFFFFF", "light", "Starter"]]
    object_headers = TABLE_SCHEMAS["Objects"][1]
    title_row = {
        "Enabled": True, "Slide_No": 1, "Object_ID": "S01_TITLE", "Object_Name": "Title", "Z_Order": 1,
        "Object_Type": "Text", "Geometry": "rect", "Text_Content": "New presentation model", "Asset_ID": None,
        "Asset_Row": None, "Left_px": 40, "Top_px": 70, "Width_px": 688, "Height_px": 100, "Rotation_deg": 0,
        "Format_ID": "TITLE", "Fill_Color": None, "Fill_Transparency": 1, "Line_Color": None,
        "Line_Transparency": 1, "Line_Width_px": 0, "Text_Alignment": "left", "Vertical_Alignment": "top",
        "Margin_Left_px": 0, "Margin_Right_px": 0, "Margin_Top_px": 0, "Margin_Bottom_px": 0,
    }
    subtitle_row = {
        **title_row,
        "Object_ID": "S01_SUBTITLE", "Object_Name": "Subtitle", "Z_Order": 2,
        "Text_Content": "Edit the model tables to define slides, objects and styles.",
        "Top_px": 185, "Height_px": 70, "Format_ID": "SUBTITLE",
    }
    object_rows = [[row.get(header) for header in object_headers] for row in [title_row, subtitle_row]]

    data_by_sheet = {
        "Control": settings_rows,
        "Formats": format_rows,
        "Slides": slide_rows,
        "Objects": object_rows,
        "Text Runs": [],
        "Paragraphs": [],
        "Assets": [],
    }
    for sheet_name, (table_name, headers) in TABLE_SCHEMAS.items():
        _write_table_sheet(workbook, sheet_name, table_name, headers, data_by_sheet[sheet_name])
    workbook.save(workbook_path(name))
    return workbook_path(name)


def import_model(model_name: str, source_workbook: Path, assets: Path | None = None, references: Path | None = None) -> Path:
    name = safe_model_name(model_name)
    source_workbook = Path(source_workbook).resolve()
    if not source_workbook.exists():
        raise FileNotFoundError(source_workbook)
    destination = model_dir(name)
    if workbook_path(name).exists():
        raise FileExistsError(f"Model already exists: {name}")
    destination.mkdir(parents=True, exist_ok=True)
    shutil.copy2(source_workbook, workbook_path(name))
    asset_dir(name).mkdir(exist_ok=True)
    reference_dir(name).mkdir(exist_ok=True)
    if assets:
        for source in Path(assets).resolve().iterdir():
            if source.is_file():
                shutil.copy2(source, asset_dir(name) / source.name)
    if references:
        for source in Path(references).resolve().iterdir():
            if source.is_file():
                shutil.copy2(source, reference_dir(name) / source.name)
    return workbook_path(name)


def clone_model(source_name: str, destination_name: str) -> Path:
    source = model_dir(source_name)
    destination = model_dir(destination_name)
    if not (source / "model.xlsx").exists():
        raise FileNotFoundError(f"Source model does not exist: {source_name}")
    if destination.exists():
        raise FileExistsError(f"Destination model already exists: {destination_name}")
    shutil.copytree(source, destination)
    return destination / "model.xlsx"


def _read_table(workbook, sheet_name: str, table_name: str) -> list[dict]:
    sheet = workbook[sheet_name]
    table = sheet.tables[table_name]
    min_col, min_row, max_col, max_row = range_boundaries(table.ref)
    headers = [sheet.cell(min_row, column).value for column in range(min_col, max_col + 1)]
    rows: list[dict] = []
    for row_number in range(min_row + 1, max_row + 1):
        values = [sheet.cell(row_number, column).value for column in range(min_col, max_col + 1)]
        if any(value is not None for value in values):
            rows.append(dict(zip(headers, values)))
    return rows


def load_model(model_name: str) -> dict:
    path = workbook_path(model_name)
    if not path.exists():
        raise FileNotFoundError(f"Model workbook not found: {path}")
    workbook = load_workbook(path, data_only=False, read_only=False)
    for sheet_name, (table_name, _) in TABLE_SCHEMAS.items():
        if sheet_name not in workbook.sheetnames:
            raise ValueError(f"Missing worksheet: {sheet_name}")
        if table_name not in workbook[sheet_name].tables:
            raise ValueError(f"Missing Excel table {table_name} on {sheet_name}")
    settings_rows = _read_table(workbook, "Control", "tblSettings")
    return {
        "settings": {str(row["Setting"]): row["Value"] for row in settings_rows},
        "formats": _read_table(workbook, "Formats", "tblFormats"),
        "slides": _read_table(workbook, "Slides", "tblSlides"),
        "objects": _read_table(workbook, "Objects", "tblObjects"),
        "runs": _read_table(workbook, "Text Runs", "tblTextRuns"),
        "paragraphs": _read_table(workbook, "Paragraphs", "tblParagraphs"),
        "assets": _read_table(workbook, "Assets", "tblAssets"),
    }


def validate_model(model_name: str) -> tuple[dict, list[str], list[str]]:
    model = load_model(model_name)
    errors: list[str] = []
    warnings: list[str] = []
    settings = model["settings"]
    width = float(settings.get("SLIDE_WIDTH_PX") or 768)
    height = float(settings.get("SLIDE_HEIGHT_PX") or 960)
    slide_numbers = [int(row["Slide_No"]) for row in model["slides"]]
    if len(slide_numbers) != len(set(slide_numbers)):
        errors.append("Slide_No values must be unique.")
    object_ids = [str(row["Object_ID"]) for row in model["objects"] if bool(row.get("Enabled"))]
    if len(object_ids) != len(set(object_ids)):
        errors.append("Enabled Object_ID values must be unique.")
    format_ids = {str(row["Format_ID"]) for row in model["formats"]}
    for row in model["objects"]:
        if not bool(row.get("Enabled")):
            continue
        object_id = str(row.get("Object_ID") or "<blank>")
        if int(row.get("Slide_No") or 0) not in slide_numbers:
            errors.append(f"{object_id}: Slide_No does not exist.")
        format_id = str(row.get("Format_ID") or "")
        if format_id and format_id not in format_ids:
            errors.append(f"{object_id}: unknown Format_ID {format_id}.")
        left = float(row.get("Left_px") or 0)
        top = float(row.get("Top_px") or 0)
        object_width = float(row.get("Width_px") or 0)
        object_height = float(row.get("Height_px") or 0)
        if left >= width or top >= height or left + object_width <= 0 or top + object_height <= 0:
            errors.append(f"{object_id}: object is completely outside the {width:g} x {height:g} canvas.")
        elif left < 0 or top < 0 or left + object_width > width + 0.01 or top + object_height > height + 0.01:
            warnings.append(f"{object_id}: geometry partially crosses the canvas edge; inspect the rendered slide.")
        if str(row.get("Object_Type") or "") in {"Image", "ChartImage"}:
            asset_name = Path(str(row.get("Asset_ID") or "")).name
            if not asset_name or not (asset_dir(model_name) / asset_name).exists():
                errors.append(f"{object_id}: missing asset {asset_name or '<blank Asset_ID>'}.")
    for row in model["runs"]:
        if str(row.get("Object_ID")) not in object_ids:
            warnings.append(f"Text run references disabled or missing object {row.get('Object_ID')}.")
    if not model["slides"]:
        errors.append("The model has no slides.")
    return model, errors, warnings


def print_status(model_name: str) -> bool:
    model, errors, warnings = validate_model(model_name)
    enabled = sum(bool(row.get("Enabled")) for row in model["objects"])
    print(f"Model: {safe_model_name(model_name)}")
    print(f"  Workbook   : {workbook_path(model_name)}")
    print(f"  Slides     : {len(model['slides'])}")
    print(f"  Objects    : {enabled}")
    print(f"  Text runs  : {len(model['runs'])}")
    print(f"  Formats    : {len(model['formats'])}")
    for warning in warnings:
        print(f"  WARNING    : {warning}")
    for error in errors:
        print(f"  ERROR      : {error}")
    print(f"  Validation : {'PASS' if not errors else 'FAIL'}")
    return not errors


def _rgb(value: str | None, fallback: str = "#000000") -> RGBColor:
    clean = str(value or fallback).replace("#", "").strip()
    if len(clean) != 6:
        clean = fallback.replace("#", "")
    return RGBColor.from_string(clean.upper())


def _paint_visible(color, transparency) -> bool:
    return bool(color) and str(color).upper() != "NONE" and float(transparency or 0) < 0.995


def _pt(px_value):
    return Pt(float(px_value or 0) * PX_TO_PT)


def _align(value):
    return {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}.get(
        str(value or "left").lower(), PP_ALIGN.LEFT
    )


def _vertical(value):
    return {"middle": MSO_ANCHOR.MIDDLE, "center": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}.get(
        str(value or "top").lower(), MSO_ANCHOR.TOP
    )


def _neutralize_theme_shadows(presentation: Presentation) -> None:
    seen_parts: set[str] = set()
    for slide_master in presentation.slide_masters:
        for relationship in slide_master.part.rels.values():
            if not relationship.reltype.endswith("/theme"):
                continue
            theme_part = relationship.target_part
            part_name = str(theme_part.partname)
            if part_name in seen_parts:
                continue
            seen_parts.add(part_name)
            root = ET.fromstring(theme_part.blob)
            changed = False
            for style_list in root.findall(f".//{{{DRAWING_NS}}}effectStyleLst"):
                for effect_style in list(style_list):
                    for child in list(effect_style):
                        effect_style.remove(child)
                    ET.SubElement(effect_style, f"{{{DRAWING_NS}}}effectLst")
                    changed = True
            if changed:
                theme_part._blob = ET.tostring(root, encoding="utf-8", xml_declaration=True)


def _style_run(run, style: dict) -> None:
    font = run.font
    font.name = str(style.get("Font_Family") or "Arial")
    font.size = Pt(float(style.get("Font_Size_pt") or 12))
    font.bold = bool(style.get("Bold"))
    font.italic = bool(style.get("Italic"))
    font.color.rgb = _rgb(style.get("Font_Color") or style.get("Color") or "#0B1826")


def _set_fill_and_line(shape, row: dict) -> None:
    if _paint_visible(row.get("Fill_Color"), row.get("Fill_Transparency")):
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(row.get("Fill_Color"))
    else:
        shape.fill.background()
    if _paint_visible(row.get("Line_Color"), row.get("Line_Transparency")) and float(row.get("Line_Width_px") or 0) > 0:
        shape.line.color.rgb = _rgb(row.get("Line_Color"))
        shape.line.width = _pt(row.get("Line_Width_px"))
    else:
        shape.line.fill.background()


def _populate_text(shape, row, slide_row, exact_runs, runs_by_object, paragraphs_by_object, formats_by_id) -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.vertical_anchor = _vertical(row.get("Vertical_Alignment"))
    text_frame.margin_left = _pt(row.get("Margin_Left_px"))
    text_frame.margin_right = _pt(row.get("Margin_Right_px"))
    text_frame.margin_top = _pt(row.get("Margin_Top_px"))
    text_frame.margin_bottom = _pt(row.get("Margin_Bottom_px"))
    object_id = str(row.get("Object_ID"))
    lines = str(row.get("Text_Content") or "").split("\n")
    dark = str(slide_row.get("Theme_Mode") or "").lower() == "dark"
    fallback = dict(formats_by_id.get(str(row.get("Format_ID")), {}))
    fallback["Font_Color"] = fallback.get("Font_Color_Dark_Page" if dark else "Font_Color_Light_Page") or "#0B1826"
    for index, line in enumerate(lines):
        paragraph_number = index + 1
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.alignment = _align(row.get("Text_Alignment"))
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        paragraph_row = paragraphs_by_object.get(object_id, {}).get(paragraph_number)
        if paragraph_row:
            paragraph.alignment = _align(paragraph_row.get("Alignment"))
            ratio = float(paragraph_row.get("Line_Spacing_Ratio") or 0)
            if ratio > 0:
                paragraph.line_spacing = ratio
        exact = runs_by_object.get(object_id, {}).get(paragraph_number, []) if exact_runs else []
        if exact:
            for run_row in exact:
                run = paragraph.add_run()
                run.text = str(run_row.get("Run_Text") or "")
                _style_run(run, run_row)
        elif line:
            run = paragraph.add_run()
            run.text = line
            _style_run(run, fallback)


def build_presentation(model_name: str, logger: Callable[[str], None] | None = None) -> Path:
    global PX_TO_PT
    model, errors, _ = validate_model(model_name)
    if errors:
        raise ValueError("Model validation failed:\n- " + "\n- ".join(errors))
    settings = model["settings"]
    PX_TO_PT = float(settings.get("PX_TO_PT") or 0.75)
    width = float(settings.get("SLIDE_WIDTH_PX") or 768)
    height = float(settings.get("SLIDE_HEIGHT_PX") or 960)
    exact_runs = bool(settings.get("USE_EXACT_RUN_FORMATTING", True))
    slides = sorted(model["slides"], key=lambda row: int(row["Slide_No"]))
    objects_by_slide = defaultdict(list)
    for row in sorted(model["objects"], key=lambda item: (int(item["Slide_No"]), int(item["Z_Order"]))):
        if bool(row.get("Enabled")):
            objects_by_slide[int(row["Slide_No"])].append(row)
    runs_by_object = defaultdict(lambda: defaultdict(list))
    for row in model["runs"]:
        runs_by_object[str(row["Object_ID"])][int(row["Paragraph_No"])].append(row)
    for paragraph_map in runs_by_object.values():
        for rows in paragraph_map.values():
            rows.sort(key=lambda row: int(row["Run_No"]))
    paragraphs_by_object = defaultdict(dict)
    for row in model["paragraphs"]:
        paragraphs_by_object[str(row["Object_ID"])][int(row["Paragraph_No"])] = row
    formats_by_id = {str(row["Format_ID"]): row for row in model["formats"]}

    presentation = Presentation()
    _neutralize_theme_shadows(presentation)
    presentation.slide_width = _pt(width)
    presentation.slide_height = _pt(height)
    blank_layout = presentation.slide_layouts[6]
    for index, slide_row in enumerate(slides, start=1):
        if logger:
            logger(f"Building slide {index} of {len(slides)}")
        slide_number = int(slide_row["Slide_No"])
        slide = presentation.slides.add_slide(blank_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = _rgb(slide_row.get("Background_Color"), "#FFFFFF")
        for row in objects_by_slide[slide_number]:
            left, top = _pt(row.get("Left_px")), _pt(row.get("Top_px"))
            object_width, object_height = _pt(row.get("Width_px")), _pt(row.get("Height_px"))
            object_type = str(row.get("Object_Type") or "Shape")
            geometry = str(row.get("Geometry") or "rect")
            if object_type in {"Image", "ChartImage"}:
                source = asset_dir(model_name) / Path(str(row.get("Asset_ID") or "")).name
                shape = slide.shapes.add_picture(str(source), left, top, object_width, object_height)
            elif object_type == "Line" or geometry == "line":
                shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left, top, left + object_width, top + object_height)
                if _paint_visible(row.get("Line_Color"), row.get("Line_Transparency")):
                    shape.line.color.rgb = _rgb(row.get("Line_Color"))
                    shape.line.width = _pt(row.get("Line_Width_px"))
                else:
                    shape.line.fill.background()
            else:
                shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if geometry == "roundRect" else MSO_SHAPE.RECTANGLE
                shape = slide.shapes.add_shape(shape_type, left, top, object_width, object_height)
                _set_fill_and_line(shape, row)
            shape.name = re.sub(r"[^a-zA-Z0-9_-]+", "_", str(row.get("Object_ID") or "Object"))
            shape.shadow.inherit = False
            shape.rotation = float(row.get("Rotation_deg") or 0)
            if str(row.get("Text_Content") or ""):
                _populate_text(shape, row, slide_row, exact_runs, runs_by_object, paragraphs_by_object, formats_by_id)

    file_name = Path(str(settings.get("OUTPUT_FILE_NAME") or f"{safe_model_name(model_name)}.pptx")).name
    if not file_name.lower().endswith(".pptx"):
        file_name += ".pptx"
    destination = OUTPUT_DIR / safe_model_name(model_name) / file_name
    destination.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(destination)
    if logger:
        logger(f"Created {destination}")
    return destination


def _preview_color(value, transparency=0, fallback="transparent") -> str:
    if not value or str(value).upper() == "NONE" or float(transparency or 0) >= 0.995:
        return fallback
    return str(value)


def generate_preview(model_name: str) -> Path:
    model, errors, _ = validate_model(model_name)
    if errors:
        raise ValueError("Model validation failed:\n- " + "\n- ".join(errors))
    slides = sorted(model["slides"], key=lambda row: int(row["Slide_No"]))
    objects_by_slide = defaultdict(list)
    for row in sorted(model["objects"], key=lambda item: (int(item["Slide_No"]), int(item["Z_Order"]))):
        if bool(row.get("Enabled")):
            objects_by_slide[int(row["Slide_No"])].append(row)
    runs_by_object = defaultdict(lambda: defaultdict(list))
    for row in model["runs"]:
        runs_by_object[str(row["Object_ID"])][int(row["Paragraph_No"])].append(row)
    formats_by_id = {str(row["Format_ID"]): row for row in model["formats"]}
    model_id = safe_model_name(model_name)
    cards = []
    for slide_row in slides:
        slide_number = int(slide_row["Slide_No"])
        elements = []
        for row in objects_by_slide[slide_number]:
            left, top = float(row.get("Left_px") or 0), float(row.get("Top_px") or 0)
            width, height_value = float(row.get("Width_px") or 0), float(row.get("Height_px") or 0)
            common = f"left:{left}px;top:{top}px;width:{width}px;height:{height_value}px;"
            object_type = str(row.get("Object_Type") or "")
            geometry = str(row.get("Geometry") or "rect")
            if object_type in {"Image", "ChartImage"}:
                asset = html.escape(Path(str(row.get("Asset_ID") or "")).name)
                elements.append(f'<img class="obj" style="{common}" src="../../models/{model_id}/assets/{asset}" alt="chart">')
            elif object_type == "Line" or geometry == "line":
                stroke = _preview_color(row.get("Line_Color"), row.get("Line_Transparency"))
                line_width = max(float(row.get("Line_Width_px") or 1), 0.5)
                elements.append(
                    f'<svg class="obj line" style="left:{left}px;top:{top}px;width:{max(width, 1)}px;height:{max(abs(height_value), 1)}px">'
                    f'<line x1="0" y1="0" x2="{width}" y2="{height_value}" stroke="{stroke}" stroke-width="{line_width}"/></svg>'
                )
            else:
                fill = _preview_color(row.get("Fill_Color"), row.get("Fill_Transparency"))
                border = _preview_color(row.get("Line_Color"), row.get("Line_Transparency"))
                border_width = float(row.get("Line_Width_px") or 0)
                radius = "10px" if geometry == "roundRect" else "0"
                elements.append(
                    f'<div class="obj" style="{common}background:{fill};border:{border_width}px solid {border};border-radius:{radius};box-sizing:border-box"></div>'
                )
            text_value = str(row.get("Text_Content") or "")
            if text_value:
                object_id = str(row.get("Object_ID"))
                dark = str(slide_row.get("Theme_Mode") or "").lower() == "dark"
                fallback = formats_by_id.get(str(row.get("Format_ID")), {})
                fallback_color = fallback.get("Font_Color_Dark_Page" if dark else "Font_Color_Light_Page") or "#0B1826"
                paragraphs = []
                for paragraph_number, line in enumerate(text_value.split("\n"), start=1):
                    runs = runs_by_object.get(object_id, {}).get(paragraph_number, [])
                    if runs:
                        spans = []
                        for run in runs:
                            style = (
                                f"font-family:{html.escape(str(run.get('Font_Family') or 'Arial'))};"
                                f"font-size:{float(run.get('Font_Size_pt') or 12)}pt;"
                                f"color:{html.escape(str(run.get('Font_Color') or '#0B1826'))};"
                                f"font-weight:{'700' if run.get('Bold') else '400'};"
                                f"font-style:{'italic' if run.get('Italic') else 'normal'};"
                            )
                            spans.append(f'<span style="{style}">{html.escape(str(run.get("Run_Text") or ""))}</span>')
                        paragraphs.append(f"<div>{''.join(spans)}</div>")
                    else:
                        style = (
                            f"font-family:{html.escape(str(fallback.get('Font_Family') or 'Arial'))};"
                            f"font-size:{float(fallback.get('Font_Size_pt') or 12)}pt;"
                            f"color:{fallback_color};font-weight:{'700' if fallback.get('Bold') else '400'};"
                        )
                        paragraphs.append(f'<div style="{style}">{html.escape(line) or "&nbsp;"}</div>')
                vertical = {"middle": "center", "bottom": "flex-end"}.get(str(row.get("Vertical_Alignment") or "top").lower(), "flex-start")
                padding = (
                    float(row.get("Margin_Top_px") or 0), float(row.get("Margin_Right_px") or 0),
                    float(row.get("Margin_Bottom_px") or 0), float(row.get("Margin_Left_px") or 0),
                )
                text_style = (
                    f"{common}display:flex;align-items:{vertical};text-align:{str(row.get('Text_Alignment') or 'left').lower()};"
                    f"padding:{padding[0]}px {padding[1]}px {padding[2]}px {padding[3]}px;box-sizing:border-box;overflow:hidden;line-height:1;"
                )
                elements.append(f'<div class="obj text" style="{text_style}"><div style="width:100%">{"".join(paragraphs)}</div></div>')
        reference = reference_dir(model_name) / f"source-slide-{slide_number:02d}.png"
        reference_html = ""
        if reference.exists():
            reference_html = (
                f'<button onclick="toggleView(this)">Show reference</button></header>'
                f'<div class="viewport model"><div class="slide" style="background:{slide_row.get("Background_Color") or "#FFFFFF"}">{"".join(elements)}</div></div>'
                f'<div class="viewport reference hidden"><img src="../../models/{model_id}/reference/{reference.name}" alt="reference slide {slide_number}"></div>'
            )
        else:
            reference_html = (
                f'</header><div class="viewport model"><div class="slide" style="background:{slide_row.get("Background_Color") or "#FFFFFF"}">{"".join(elements)}</div></div>'
            )
        cards.append(f'<section class="card"><header><strong>Slide {slide_number}</strong>{reference_html}</section>')
    page = f"""<!doctype html><html lang="en"><head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>{html.escape(model_id)} preview</title><style>
body{{margin:0;background:#eef2f5;color:#0b1826;font:14px Arial,sans-serif}}.top{{position:sticky;top:0;z-index:9;background:#071b27;color:white;padding:16px 24px}}.top h1{{margin:0;font-size:20px}}main{{padding:24px;display:grid;grid-template-columns:repeat(auto-fit,minmax(410px,1fr));gap:24px}}.card{{background:white;border:1px solid #d7e0e7;border-radius:8px;padding:14px}}header{{display:flex;justify-content:space-between;align-items:center;margin-bottom:12px}}button{{border:1px solid #2f6bff;background:white;color:#2f6bff;padding:7px 10px;border-radius:4px}}.viewport{{width:384px;height:480px;margin:auto;overflow:hidden;background:#ddd}}.slide{{position:relative;width:768px;height:960px;transform:scale(.5);transform-origin:top left;overflow:hidden}}.obj{{position:absolute;margin:0;padding:0}}.line{{overflow:visible}}.reference img{{display:block;width:384px;height:480px}}.hidden{{display:none}}
</style></head><body><div class="top"><h1>{html.escape(model_id)} — model preview</h1></div><main>{''.join(cards)}</main>
<script>function toggleView(button){{const card=button.closest('.card');const model=card.querySelector('.model');const ref=card.querySelector('.reference');model.classList.toggle('hidden');ref.classList.toggle('hidden');button.textContent=model.classList.contains('hidden')?'Show model':'Show reference';}}</script></body></html>"""
    destination = PREVIEW_DIR / model_id / "index.html"
    destination.parent.mkdir(parents=True, exist_ok=True)
    destination.write_text(page, encoding="utf-8")
    return destination


def open_path(path: Path) -> None:
    path = path.resolve()
    if os.name == "nt":
        os.startfile(path)  # type: ignore[attr-defined]
    else:
        raise RuntimeError(f"Open manually: {path}")


def interactive_menu() -> None:
    ensure_root_structure()
    while True:
        models = list_models()
        print("\nPresentation Model Manager")
        print("  Models: " + (", ".join(models) if models else "none"))
        print("  1. Create model")
        print("  2. Clone model")
        print("  3. Validate model")
        print("  4. Build presentation")
        print("  5. Generate preview")
        print("  6. Open model workbook")
        print("  7. Open GUI manager")
        print("  0. Exit")
        choice = input("Select an option: ").strip()
        try:
            if choice == "0":
                return
            if choice == "1":
                print(f"Created: {create_model(input('New model name: '))}")
            elif choice == "2":
                print(f"Cloned: {clone_model(input('Source model: '), input('New model name: '))}")
            elif choice == "3":
                print_status(input(f"Model [{DEFAULT_MODEL}]: ").strip() or DEFAULT_MODEL)
            elif choice == "4":
                print(build_presentation(input(f"Model [{DEFAULT_MODEL}]: ").strip() or DEFAULT_MODEL, print))
            elif choice == "5":
                name = input(f"Model [{DEFAULT_MODEL}]: ").strip() or DEFAULT_MODEL
                destination = generate_preview(name)
                print(destination)
            elif choice == "6":
                open_path(workbook_path(input(f"Model [{DEFAULT_MODEL}]: ").strip() or DEFAULT_MODEL))
            elif choice == "7":
                launch_gui()
            else:
                print("Invalid option.")
        except Exception as exc:
            print(f"ERROR: {exc}", file=sys.stderr)


def launch_gui() -> None:
    try:
        import tkinter as tk
        from tkinter import messagebox, simpledialog, ttk
    except ModuleNotFoundError as exc:
        raise RuntimeError("Tkinter is not available in this Python environment.") from exc

    class Manager(tk.Tk):
        def __init__(self):
            super().__init__()
            self.title("Presentation Model Manager")
            self.geometry("860x610")
            self.configure(bg="#EEF2F5")
            header = tk.Frame(self, bg="#071B27", padx=24, pady=18)
            header.pack(fill="x")
            tk.Label(header, text="Presentation Model Manager", bg="#071B27", fg="white", font=("Arial", 21, "bold")).pack(anchor="w")
            tk.Label(header, text="One Python file · spreadsheet models · no BAT or ActiveX", bg="#071B27", fg="#B9C8D6", font=("Arial", 10)).pack(anchor="w")
            body = tk.Frame(self, bg="#EEF2F5", padx=24, pady=20)
            body.pack(fill="both", expand=True)
            selection = tk.Frame(body, bg="#EEF2F5")
            selection.pack(fill="x")
            tk.Label(selection, text="Model", bg="#EEF2F5", font=("Arial", 10, "bold")).pack(side="left")
            self.model_value = tk.StringVar()
            self.model_box = ttk.Combobox(selection, textvariable=self.model_value, state="readonly", width=38)
            self.model_box.pack(side="left", padx=10)
            ttk.Button(selection, text="Refresh", command=self.refresh_models).pack(side="left")
            buttons = tk.Frame(body, bg="#EEF2F5")
            buttons.pack(fill="x", pady=16)
            actions = [
                ("Create model", self.create_action), ("Clone model", self.clone_action),
                ("Validate", self.validate_action), ("Build PPTX", self.build_action),
                ("Preview", self.preview_action), ("Open workbook", self.open_workbook_action),
            ]
            for index, (label, command) in enumerate(actions):
                ttk.Button(buttons, text=label, command=command).grid(row=index // 3, column=index % 3, padx=5, pady=5, sticky="ew")
            for column in range(3):
                buttons.columnconfigure(column, weight=1)
            self.log_box = tk.Text(body, height=20, bg="white", fg="#0B1826", relief="flat", padx=12, pady=10, font=("Consolas", 9), state="disabled")
            self.log_box.pack(fill="both", expand=True)
            self.refresh_models()

        def log(self, message: str):
            self.log_box.configure(state="normal")
            self.log_box.insert("end", message + "\n")
            self.log_box.see("end")
            self.log_box.configure(state="disabled")

        def selected(self) -> str:
            value = self.model_value.get().strip()
            if not value:
                raise ValueError("Select a model first.")
            return value

        def refresh_models(self):
            values = list_models()
            self.model_box["values"] = values
            if values and self.model_value.get() not in values:
                self.model_value.set(DEFAULT_MODEL if DEFAULT_MODEL in values else values[0])

        def create_action(self):
            value = simpledialog.askstring("Create model", "Model name:", parent=self)
            if value:
                try:
                    self.log(f"Created {create_model(value)}")
                    self.refresh_models()
                except Exception as exc:
                    messagebox.showerror("Create model", str(exc), parent=self)

        def clone_action(self):
            try:
                source = self.selected()
                value = simpledialog.askstring("Clone model", "New model name:", parent=self)
                if value:
                    self.log(f"Cloned {clone_model(source, value)}")
                    self.refresh_models()
            except Exception as exc:
                messagebox.showerror("Clone model", str(exc), parent=self)

        def validate_action(self):
            try:
                self.log("Validation PASS" if print_status(self.selected()) else "Validation FAIL")
            except Exception as exc:
                messagebox.showerror("Validate", str(exc), parent=self)

        def build_action(self):
            def work():
                try:
                    output = build_presentation(self.selected(), lambda message: self.after(0, self.log, message))
                    self.after(0, self.log, f"Created {output}")
                except Exception as exc:
                    self.after(0, messagebox.showerror, "Build", str(exc))
            threading.Thread(target=work, daemon=True).start()

        def preview_action(self):
            try:
                output = generate_preview(self.selected())
                self.log(f"Created {output}")
                webbrowser.open(output.as_uri())
            except Exception as exc:
                messagebox.showerror("Preview", str(exc), parent=self)

        def open_workbook_action(self):
            try:
                open_path(workbook_path(self.selected()))
            except Exception as exc:
                messagebox.showerror("Open workbook", str(exc), parent=self)

    Manager().mainloop()


def main() -> None:
    ensure_root_structure()
    parser = argparse.ArgumentParser(description="Create, manage, validate and build spreadsheet-driven presentation models.")
    subparsers = parser.add_subparsers(dest="command")
    subparsers.add_parser("list", help="List registered models.")
    create_parser = subparsers.add_parser("create", help="Create a new model and workbook structure.")
    create_parser.add_argument("name")
    import_parser = subparsers.add_parser("import", help="Register an existing model workbook.")
    import_parser.add_argument("name")
    import_parser.add_argument("workbook", type=Path)
    import_parser.add_argument("--assets", type=Path)
    import_parser.add_argument("--references", type=Path)
    clone_parser = subparsers.add_parser("clone", help="Clone a model profile.")
    clone_parser.add_argument("source")
    clone_parser.add_argument("destination")
    status_parser = subparsers.add_parser("status", help="Show counts and validate a model.")
    status_parser.add_argument("name", nargs="?", default=DEFAULT_MODEL)
    validate_parser = subparsers.add_parser("validate", help="Validate a model.")
    validate_parser.add_argument("name", nargs="?", default=DEFAULT_MODEL)
    build_parser = subparsers.add_parser("build", help="Build a model into a PPTX file.")
    build_parser.add_argument("name", nargs="?", default=DEFAULT_MODEL)
    preview_parser = subparsers.add_parser("preview", help="Generate a browser preview.")
    preview_parser.add_argument("name", nargs="?", default=DEFAULT_MODEL)
    preview_parser.add_argument("--open", action="store_true")
    open_parser = subparsers.add_parser("open", help="Open a model workbook in Windows.")
    open_parser.add_argument("name", nargs="?", default=DEFAULT_MODEL)
    subparsers.add_parser("gui", help="Open the optional graphical manager.")
    args = parser.parse_args()
    try:
        if args.command is None:
            interactive_menu()
        elif args.command == "list":
            for name in list_models():
                print(name)
        elif args.command == "create":
            print(create_model(args.name))
        elif args.command == "import":
            print(import_model(args.name, args.workbook, args.assets, args.references))
        elif args.command == "clone":
            print(clone_model(args.source, args.destination))
        elif args.command in {"status", "validate"}:
            if not print_status(args.name):
                raise SystemExit(1)
        elif args.command == "build":
            print(build_presentation(args.name, print))
        elif args.command == "preview":
            output = generate_preview(args.name)
            print(output)
            if args.open:
                webbrowser.open(output.as_uri())
        elif args.command == "open":
            open_path(workbook_path(args.name))
        elif args.command == "gui":
            launch_gui()
    except Exception as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        raise SystemExit(1) from exc


if __name__ == "__main__":
    main()
