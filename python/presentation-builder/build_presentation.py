from __future__ import annotations

import argparse
from collections import defaultdict
from pathlib import Path
from typing import Callable, Iterable
from xml.etree import ElementTree as ET

from openpyxl import load_workbook
from openpyxl.utils.cell import range_boundaries
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Pt


ROOT = Path(__file__).resolve().parent
DEFAULT_CONFIG = ROOT / "config" / "presentation_config.xlsx"
DEFAULT_OUTPUT = ROOT / "output" / "Generated_USD_Cash_Allocation.pptx"
ASSET_DIR = ROOT / "assets"
PX_TO_PT = 0.75
DRAWING_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _log(callback: Callable[[str], None] | None, message: str) -> None:
    if callback:
        callback(message)
    else:
        print(message)


def _read_table(workbook, sheet_name: str, table_name: str) -> list[dict]:
    sheet = workbook[sheet_name]
    table = sheet.tables[table_name]
    min_col, min_row, max_col, max_row = range_boundaries(table.ref)
    headers = [sheet.cell(min_row, col).value for col in range(min_col, max_col + 1)]
    rows: list[dict] = []
    for row_no in range(min_row + 1, max_row + 1):
        values = [sheet.cell(row_no, col).value for col in range(min_col, max_col + 1)]
        if any(value is not None for value in values):
            rows.append(dict(zip(headers, values)))
    return rows


def load_model(config_path: Path = DEFAULT_CONFIG) -> dict:
    workbook = load_workbook(config_path, data_only=False, read_only=False)
    settings_rows = _read_table(workbook, "Control", "tblSettings")
    settings = {str(row["Setting"]): row["Value"] for row in settings_rows}
    return {
        "settings": settings,
        "slides": _read_table(workbook, "Slides", "tblSlides"),
        "objects": _read_table(workbook, "Objects", "tblObjects"),
        "runs": _read_table(workbook, "Text Runs", "tblTextRuns"),
        "paragraphs": _read_table(workbook, "Paragraphs", "tblParagraphs"),
        "formats": _read_table(workbook, "Formats", "tblFormats"),
    }


def configured_output_path(config_path: Path = DEFAULT_CONFIG) -> Path:
    settings = load_model(config_path)["settings"]
    file_name = Path(str(settings.get("OUTPUT_FILE_NAME") or DEFAULT_OUTPUT.name)).name
    if file_name.lower().endswith(".pptx") is False:
        file_name += ".pptx"
    return ROOT / "output" / file_name


def _rgb(value: str | None, fallback: str = "#000000") -> RGBColor:
    clean = str(value or fallback).replace("#", "").strip()
    if len(clean) != 6:
        clean = fallback.replace("#", "")
    return RGBColor.from_string(clean.upper())


def _paint_is_visible(color, transparency) -> bool:
    return bool(color) and str(color).upper() != "NONE" and float(transparency or 0) < 0.995


def _pt(px_value) -> int:
    return Pt(float(px_value or 0) * PX_TO_PT)


def _align(value) -> PP_ALIGN:
    return {
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }.get(str(value or "left").lower(), PP_ALIGN.LEFT)


def _vertical(value) -> MSO_ANCHOR:
    return {
        "middle": MSO_ANCHOR.MIDDLE,
        "center": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get(str(value or "top").lower(), MSO_ANCHOR.TOP)


def _safe_name(value) -> str:
    return str(value or "Object").replace("/", "_").replace(" ", "_")


def _neutralize_theme_shadows(presentation: Presentation) -> None:
    """Remove inherited Office-theme shadows from newly created objects."""
    seen_parts: set[str] = set()
    for slide_master in presentation.slide_masters:
        for relationship in slide_master.part.rels.values():
            if not relationship.reltype.endswith("/theme"):
                continue
            theme_part = relationship.target_part
            part_name = str(theme_part.partname)
            if part_name in seen_parts:
                continue
            seen_parts.add(part_name)
            root = ET.fromstring(theme_part.blob)
            changed = False
            for style_list in root.findall(f".//{{{DRAWING_NS}}}effectStyleLst"):
                for effect_style in list(style_list):
                    for child in list(effect_style):
                        effect_style.remove(child)
                    ET.SubElement(effect_style, f"{{{DRAWING_NS}}}effectLst")
                    changed = True
            if changed:
                theme_part._blob = ET.tostring(root, encoding="utf-8", xml_declaration=True)


def _style_run(run, style: dict) -> None:
    font = run.font
    font.name = str(style.get("Font_Family") or "Aptos")
    font.size = Pt(float(style.get("Font_Size_pt") or 12))
    font.bold = bool(style.get("Bold"))
    font.italic = bool(style.get("Italic"))
    color = style.get("Font_Color") or style.get("Color") or "#0B1826"
    font.color.rgb = _rgb(str(color))


def _set_fill_and_line(shape, row: dict) -> None:
    if _paint_is_visible(row.get("Fill_Color"), row.get("Fill_Transparency")):
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(row.get("Fill_Color"))
    else:
        shape.fill.background()

    if _paint_is_visible(row.get("Line_Color"), row.get("Line_Transparency")) and float(row.get("Line_Width_px") or 0) > 0:
        shape.line.color.rgb = _rgb(row.get("Line_Color"))
        shape.line.width = _pt(row.get("Line_Width_px"))
    else:
        shape.line.fill.background()


def _format_for_object(row: dict, slide_row: dict, format_by_id: dict[str, dict]) -> dict:
    style = dict(format_by_id.get(str(row.get("Format_ID")), {}))
    dark = str(slide_row.get("Theme_Mode") or "").lower() == "dark"
    style["Font_Color"] = style.get("Font_Color_Dark_Page" if dark else "Font_Color_Light_Page") or "#0B1826"
    return style


def _populate_text(
    shape,
    row: dict,
    slide_row: dict,
    exact_runs: bool,
    runs_by_object: dict[str, dict[int, list[dict]]],
    paragraphs_by_object: dict[str, dict[int, dict]],
    format_by_id: dict[str, dict],
) -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.vertical_anchor = _vertical(row.get("Vertical_Alignment"))
    text_frame.margin_left = _pt(row.get("Margin_Left_px"))
    text_frame.margin_right = _pt(row.get("Margin_Right_px"))
    text_frame.margin_top = _pt(row.get("Margin_Top_px"))
    text_frame.margin_bottom = _pt(row.get("Margin_Bottom_px"))

    object_id = str(row.get("Object_ID"))
    lines = str(row.get("Text_Content") or "").split("\n")
    fallback_style = _format_for_object(row, slide_row, format_by_id)

    for index, line in enumerate(lines):
        paragraph_no = index + 1
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.alignment = _align(row.get("Text_Alignment"))
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        paragraph_row = paragraphs_by_object.get(object_id, {}).get(paragraph_no)
        if paragraph_row:
            paragraph.alignment = _align(paragraph_row.get("Alignment"))
            ratio = float(paragraph_row.get("Line_Spacing_Ratio") or 0)
            if ratio > 0:
                paragraph.line_spacing = ratio

        exact = runs_by_object.get(object_id, {}).get(paragraph_no, []) if exact_runs else []
        if exact:
            for run_row in exact:
                run = paragraph.add_run()
                run.text = str(run_row.get("Run_Text") or "")
                _style_run(run, run_row)
        elif line:
            run = paragraph.add_run()
            run.text = line
            _style_run(run, fallback_style)


def build_presentation(
    config_path: Path = DEFAULT_CONFIG,
    output_path: Path | None = None,
    logger: Callable[[str], None] | None = None,
) -> Path:
    global PX_TO_PT
    config_path = Path(config_path).resolve()
    _log(logger, f"Reading {config_path.name}")
    model = load_model(config_path)
    settings = model["settings"]
    PX_TO_PT = float(settings.get("PX_TO_PT") or 0.75)
    if output_path is None:
        file_name = Path(str(settings.get("OUTPUT_FILE_NAME") or DEFAULT_OUTPUT.name)).name
        if file_name.lower().endswith(".pptx") is False:
            file_name += ".pptx"
        output_path = ROOT / "output" / file_name
    output_path = Path(output_path).resolve()
    exact_runs = bool(settings.get("USE_EXACT_RUN_FORMATTING", True))
    width_px = float(settings.get("SLIDE_WIDTH_PX", 768))
    height_px = float(settings.get("SLIDE_HEIGHT_PX", 960))

    slide_rows = sorted(model["slides"], key=lambda item: int(item["Slide_No"]))
    object_rows = sorted(
        [row for row in model["objects"] if bool(row.get("Enabled"))],
        key=lambda item: (int(item["Slide_No"]), int(item["Z_Order"])),
    )
    objects_by_slide: dict[int, list[dict]] = defaultdict(list)
    for row in object_rows:
        objects_by_slide[int(row["Slide_No"])].append(row)

    runs_by_object: dict[str, dict[int, list[dict]]] = defaultdict(lambda: defaultdict(list))
    for row in model["runs"]:
        runs_by_object[str(row["Object_ID"])][int(row["Paragraph_No"])].append(row)
    for paragraph_map in runs_by_object.values():
        for rows in paragraph_map.values():
            rows.sort(key=lambda item: int(item["Run_No"]))

    paragraphs_by_object: dict[str, dict[int, dict]] = defaultdict(dict)
    for row in model["paragraphs"]:
        paragraphs_by_object[str(row["Object_ID"])][int(row["Paragraph_No"])] = row
    format_by_id = {str(row["Format_ID"]): row for row in model["formats"]}

    presentation = Presentation()
    _neutralize_theme_shadows(presentation)
    presentation.slide_width = _pt(width_px)
    presentation.slide_height = _pt(height_px)
    blank_layout = presentation.slide_layouts[6]

    for slide_row in slide_rows:
        slide_no = int(slide_row["Slide_No"])
        _log(logger, f"Building slide {slide_no} of {len(slide_rows)}")
        slide = presentation.slides.add_slide(blank_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = _rgb(slide_row.get("Background_Color"), "#071B27")

        for row in objects_by_slide[slide_no]:
            left, top = _pt(row.get("Left_px")), _pt(row.get("Top_px"))
            width, height = _pt(row.get("Width_px")), _pt(row.get("Height_px"))
            object_type = str(row.get("Object_Type") or "Shape")
            geometry = str(row.get("Geometry") or "rect")

            if object_type in {"Image", "ChartImage"}:
                asset_name = Path(str(row.get("Asset_ID") or "")).name
                asset_path = ASSET_DIR / asset_name
                if not asset_path.exists():
                    raise FileNotFoundError(f"Missing embedded chart asset: {asset_path}")
                shape = slide.shapes.add_picture(str(asset_path), left, top, width, height)
            elif object_type == "Line" or geometry == "line":
                shape = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    left,
                    top,
                    left + width,
                    top + height,
                )
                if _paint_is_visible(row.get("Line_Color"), row.get("Line_Transparency")):
                    shape.line.color.rgb = _rgb(row.get("Line_Color"))
                    shape.line.width = _pt(row.get("Line_Width_px"))
                else:
                    shape.line.fill.background()
            else:
                shape_kind = MSO_SHAPE.ROUNDED_RECTANGLE if geometry == "roundRect" else MSO_SHAPE.RECTANGLE
                shape = slide.shapes.add_shape(shape_kind, left, top, width, height)
                _set_fill_and_line(shape, row)

            shape.name = _safe_name(row.get("Object_ID"))
            shape.shadow.inherit = False
            shape.rotation = float(row.get("Rotation_deg") or 0)
            if str(row.get("Text_Content") or ""):
                _populate_text(
                    shape,
                    row,
                    slide_row,
                    exact_runs,
                    runs_by_object,
                    paragraphs_by_object,
                    format_by_id,
                )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)
    _log(logger, f"Created {output_path.name}")
    return output_path


def main() -> None:
    parser = argparse.ArgumentParser(description="Build the spreadsheet-driven PowerPoint without ActiveX or PowerPoint.")
    parser.add_argument("--config", type=Path, default=DEFAULT_CONFIG)
    parser.add_argument("--output", type=Path)
    args = parser.parse_args()
    build_presentation(args.config, args.output)


if __name__ == "__main__":
    main()
